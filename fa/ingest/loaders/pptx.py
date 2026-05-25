"""PowerPoint 抽文 — python-pptx."""

from pathlib import Path


def load_pptx(path: Path) -> tuple[str, int]:
    """返回 (纯文本, slide 数)。每页拼接标题 + 正文 + 备注。"""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    n = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {i}"]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    slide_parts.append(text)
        # 备注页
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"\n_(备注)_: {notes}")
        parts.append("\n".join(slide_parts))
    return "\n\n".join(parts), n
